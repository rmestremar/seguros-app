"""
Mejora visual - Presentación Santa Lucía / Caso 2 Javier
================================================================
- Transiciones Fade suaves en todas las diapositivas
- Animaciones de entrada auto-play escalonadas
- Slide 3: animación dramática del ascenso de Javier (top-right)
  con estrellas decorativas y secuencia ANTES→ASCENSO→AHORA
- Slide 4: mejora de tabla (tipografía, alturas de fila, layout)
- NO modifica el contenido de texto
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from lxml import etree

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ── ID counter global ────────────────────────────────────────────────────────
_id_counter = [200]  # start alto para evitar colisión con IDs existentes

def next_id():
    _id_counter[0] += 1
    return str(_id_counter[0])

# ════════════════════════════════════════════════════════════════════════════
# TRANSICIONES
# ════════════════════════════════════════════════════════════════════════════
def add_fade_transition(slide, speed="fast"):
    xml = f'<p:transition xmlns:p="{P_NS}" spd="{speed}"><p:fade/></p:transition>'
    slide._element.append(etree.fromstring(xml))


# ════════════════════════════════════════════════════════════════════════════
# BLOQUES DE ANIMACIÓN
# ════════════════════════════════════════════════════════════════════════════
def _fade_block(shape_id, delay_ms, dur_ms=500, grp=0, node_type="afterEffect"):
    id1, id2, id3, id4 = next_id(), next_id(), next_id(), next_id()
    return etree.fromstring(f"""<p:par xmlns:p="{P_NS}">
  <p:cTn id="{id1}" fill="hold">
    <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{id2}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{id3}" presetID="10" presetClass="entr" presetSubtype="0"
                     fill="hold" grpId="{grp}" nodeType="{node_type}">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:animEffect transition="in" filter="fade">
                    <p:cBhvr>
                      <p:cTn id="{id4}" dur="{dur_ms}" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                    </p:cBhvr>
                  </p:animEffect>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>""")


def _fly_bottom_block(shape_id, delay_ms, dur_ms=600, grp=0, node_type="afterEffect"):
    """Fly in desde abajo hacia arriba (representa ascenso)."""
    id1, id2, id3, id4, id5, id6 = [next_id() for _ in range(6)]
    return etree.fromstring(f"""<p:par xmlns:p="{P_NS}">
  <p:cTn id="{id1}" fill="hold">
    <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{id2}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{id3}" presetID="2" presetClass="entr" presetSubtype="8"
                     fill="hold" grpId="{grp}" nodeType="{node_type}">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="{id4}" dur="1" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:animMotion origin="parent" path="M 0 0.25 L 0 0" pathEditMode="relative" dir="0">
                    <p:cBhvr>
                      <p:cTn id="{id5}" dur="{dur_ms}" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                      <p:attrNameLst>
                        <p:attrName>ppt_x</p:attrName>
                        <p:attrName>ppt_y</p:attrName>
                      </p:attrNameLst>
                    </p:cBhvr>
                  </p:animMotion>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>""")


def _fly_left_block(shape_id, delay_ms, dur_ms=600, grp=0, node_type="afterEffect"):
    """Fly in desde la izquierda."""
    id1, id2, id3, id4, id5 = [next_id() for _ in range(5)]
    return etree.fromstring(f"""<p:par xmlns:p="{P_NS}">
  <p:cTn id="{id1}" fill="hold">
    <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{id2}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{id3}" presetID="2" presetClass="entr" presetSubtype="2"
                     fill="hold" grpId="{grp}" nodeType="{node_type}">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="{id4}" dur="1" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:animMotion origin="parent" path="M -0.3 0 L 0 0" pathEditMode="relative" dir="0">
                    <p:cBhvr>
                      <p:cTn id="{id5}" dur="{dur_ms}" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                      <p:attrNameLst>
                        <p:attrName>ppt_x</p:attrName>
                        <p:attrName>ppt_y</p:attrName>
                      </p:attrNameLst>
                    </p:cBhvr>
                  </p:animMotion>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>""")


def build_timing(effects):
    """Envuelve la lista de efectos en el bloque <p:timing> completo."""
    id_root, id_seq = next_id(), next_id()
    xml = f"""<p:timing xmlns:p="{P_NS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="{id_root}" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="{id_seq}" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst/>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrevClick" delay="0"><p:tn/></p:cond>
            </p:prevCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst/>
</p:timing>"""
    timing_el = etree.fromstring(xml)
    ctnlst = timing_el.find(f".//{{{P_NS}}}seq/{{{P_NS}}}cTn/{{{P_NS}}}childTnLst")
    for eff in effects:
        ctnlst.append(eff)
    return timing_el


def add_animations(slide, specs):
    """
    specs: list of (shape_id, delay_ms, effect)
    effect: 'fade' | 'fly_bottom' | 'fly_left'
    """
    effects = []
    for i, spec in enumerate(specs):
        sid, delay = spec[0], spec[1]
        eff_type = spec[2] if len(spec) > 2 else "fade"
        dur = spec[3] if len(spec) > 3 else (500 if eff_type == "fade" else 600)
        grp = i

        if eff_type == "fly_bottom":
            effects.append(_fly_bottom_block(sid, delay, dur, grp))
        elif eff_type == "fly_left":
            effects.append(_fly_left_block(sid, delay, dur, grp))
        else:
            effects.append(_fade_block(sid, delay, dur, grp))

    slide._element.append(build_timing(effects))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Animación del ascenso de Javier
# ════════════════════════════════════════════════════════════════════════════
def _get_max_shape_id(slide):
    ids = []
    for el in slide._element.iter():
        id_val = el.get("id")
        if id_val and id_val.isdigit():
            ids.append(int(id_val))
    return max(ids) if ids else 100


def add_star_shape(slide, x_emu, y_emu, size_emu, color_hex, opacity=100000):
    """Añade una estrella de 5 puntas decorativa en la diapositiva."""
    new_id = _get_max_shape_id(slide) + 1
    sp_tree = slide.shapes._spTree
    xml = f"""<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="{new_id}" name="StarDeco_{new_id}"/>
    <p:cNvSpPr/><p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x_emu}" y="{y_emu}"/><a:ext cx="{size_emu}" cy="{size_emu}"/></a:xfrm>
    <a:prstGeom prst="star5"><a:avLst/></a:prstGeom>
    <a:solidFill>
      <a:srgbClr val="{color_hex}">
        <a:alpha val="{opacity}"/>
      </a:srgbClr>
    </a:solidFill>
    <a:ln><a:noFill/></a:ln>
    <a:effectLst>
      <a:outerShdw blurRad="38100" dist="12700" dir="5400000" rotWithShape="0">
        <a:srgbClr val="{color_hex}"><a:alpha val="30000"/></a:srgbClr>
      </a:outerShdw>
    </a:effectLst>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>"""
    el = etree.fromstring(xml)
    sp_tree.append(el)
    return new_id


def add_glow_rect(slide, x_emu, y_emu, w_emu, h_emu, color_hex):
    """Añade un rectángulo de resplandor detrás de la tarjeta AHORA."""
    new_id = _get_max_shape_id(slide) + 1
    sp_tree = slide.shapes._spTree
    xml = f"""<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="{new_id}" name="GlowRect_{new_id}"/>
    <p:cNvSpPr/><p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x_emu}" y="{y_emu}"/><a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 6000"/></a:avLst></a:prstGeom>
    <a:solidFill>
      <a:srgbClr val="{color_hex}"><a:alpha val="12000"/></a:srgbClr>
    </a:solidFill>
    <a:ln><a:noFill/></a:ln>
    <a:effectLst>
      <a:outerShdw blurRad="228600" dist="0" dir="0" rotWithShape="0">
        <a:srgbClr val="{color_hex}"><a:alpha val="25000"/></a:srgbClr>
      </a:outerShdw>
    </a:effectLst>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>"""
    el = etree.fromstring(xml)
    # Insertar ANTES del resto de shapes para que quede por detrás
    sp_tree.insert(2, el)
    return new_id


def add_trajectory_arrow(slide, x1, y1, x2, y2, color_hex):
    """Añade una línea con flecha diagonal entre dos puntos."""
    new_id = _get_max_shape_id(slide) + 1
    sp_tree = slide.shapes._spTree
    # flip horizontal/vertical según dirección
    flip = ""
    cx = abs(x2 - x1)
    cy = abs(y2 - y1)
    xml = f"""<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
  <p:nvSpPr>
    <p:cNvPr id="{new_id}" name="Arrow_{new_id}"/>
    <p:cNvSpPr/><p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x1}" y="{y1}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="28575" cap="rnd">
      <a:solidFill>
        <a:srgbClr val="{color_hex}"><a:alpha val="55000"/></a:srgbClr>
      </a:solidFill>
      <a:prstDash val="sysDot"/>
      <a:tailEnd type="arrow" w="med" len="med"/>
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>"""
    el = etree.fromstring(xml)
    sp_tree.append(el)
    return new_id


def enhance_slide3_javier(slide):
    """
    Crea la animación dramática del ascenso en el área top-right.
    Grupos existentes:
      id=40  Grupo 39 → ANTES  (Javier · Vendedor laboral)
      id=41  Grupo 40 → ASCENSO badge
      id=42  Grupo 41 → AHORA  (Javier · Responsable del equipo)
    """
    # Posiciones de referencia (EMU)
    ANTES_L, ANTES_T = 5985405, 157014
    ANTES_W, ANTES_H = 4206240, 1417320
    AHORA_L, AHORA_T = 8086346, 1192060
    AHORA_W, AHORA_H = 4206240, 1417320

    # 1. Rectángulo de resplandor azul detrás de AHORA (se añade primero = detrás)
    glow_id = add_glow_rect(slide,
                            AHORA_L - 76200, AHORA_T - 76200,
                            AHORA_W + 152400, AHORA_H + 152400,
                            "0055A4")

    # 2. Flecha de trayectoria diagonal ANTES→AHORA
    arrow_id = add_trajectory_arrow(slide,
                                    ANTES_L + ANTES_W - 200000,  # parte derecha de ANTES
                                    ANTES_T + ANTES_H // 2,       # centro vertical
                                    AHORA_L + 200000,              # parte izquierda de AHORA
                                    AHORA_T + AHORA_H // 2,
                                    "2E9BE6")

    # 3. Estrellas decorativas (5-pointed, gold) alrededor de AHORA
    star_gold  = "E07B00"
    star_blue  = "2E9BE6"
    star_navy  = "0055A4"

    # cluster de estrellas en la esquina superior derecha cerca de AHORA
    s1 = add_star_shape(slide, 10350000, 80000,  350000, star_gold,  100000)
    s2 = add_star_shape(slide, 10850000, 220000, 230000, star_gold,  85000)
    s3 = add_star_shape(slide, 11200000, 50000,  180000, star_blue,  90000)
    s4 = add_star_shape(slide, 11550000, 280000, 140000, star_gold,  75000)
    s5 = add_star_shape(slide, 10600000, 450000, 150000, star_navy,  70000)

    # ── Secuencia de animación ──────────────────────────────────────────────
    specs = [
        # ANTES: fade suave (representa el pasado)
        (40,    0,    "fade",       700),
        # flecha aparece después de ANTES
        (arrow_id, 600, "fade",    500),
        # ASCENSO sube dramáticamente desde abajo
        (41,    900,  "fly_bottom", 500),
        # resplandor detrás de AHORA aparece
        (glow_id, 1400, "fade",    400),
        # AHORA entra con fly desde la derecha (representando la llegada)
        (42,    1400, "fly_left",   600),
        # estrellas aparecen en cascada
        (s1,    2100, "fade",      300),
        (s2,    2250, "fade",      300),
        (s3,    2400, "fade",      300),
        (s4,    2550, "fade",      300),
        (s5,    2700, "fade",      300),
    ]
    add_animations(slide, specs)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Mejora de la tabla
# ════════════════════════════════════════════════════════════════════════════
def improve_slide4_table(slide):
    """Mejora tipografía y alturas de la tabla de datos."""
    for shape in slide.shapes:
        if shape.shape_type != 19:  # TABLE
            continue

        table = shape.table

        # ── Actualizar alturas de filas ─────────────────────────────────────
        # Row 0 (sub-headers): +25%
        # Row 1 (col headers): +30%
        # Rows 2-5 (datos):    +35%
        # Row 6 (EQUIPO):      +30%
        row_scale = [1.25, 1.30, 1.35, 1.35, 1.35, 1.35, 1.30]
        total_new_h = 0

        for i, (row, scale) in enumerate(zip(table.rows, row_scale)):
            tr = row._tr
            old_h = int(tr.get("h", 0))
            new_h = int(old_h * scale)
            tr.set("h", str(new_h))
            total_new_h += new_h

        # Actualizar altura total del graphicFrame
        xfrm = shape._element.find(f"{{{P_NS}}}xfrm")
        if xfrm is None:
            xfrm = shape._element.find(f".//{{{A_NS}}}off/..")
        # Buscar el ext en el graphicFrame directamente
        for el in shape._element:
            if el.tag.endswith("}xfrm"):
                for sub in el:
                    if sub.tag.endswith("}ext"):
                        sub.set("cy", str(total_new_h))
                break

        # ── Actualizar tamaños de fuente (directo en XML) ───────────────────
        # Mapeo hundredths-of-pt → nuevos hundredths-of-pt
        size_map = {
            "800":  "1000",   # 8pt   → 10pt
            "850":  "1100",   # 8.5pt → 11pt  (headers de columna)
            "1000": "1200",   # 10pt  → 12pt  (puntos)
            "1100": "1400",   # 11pt  → 14pt  (datos numéricos)
            "1200": "1500",   # 12pt  → 15pt  (nombres y conversiones)
        }

        for row in table.rows:
            for cell in row.cells:
                for el in cell._tc.iter():
                    sz = el.get("sz")
                    if sz and sz in size_map:
                        el.set("sz", size_map[sz])

        # ── Mejorar padding de celdas para mayor legibilidad ────────────────
        for row in table.rows:
            for cell in row.cells:
                tc = cell._tc
                tcPr = tc.find(f"{{{A_NS}}}tcPr")
                if tcPr is None:
                    tcPr = etree.SubElement(tc, f"{{{A_NS}}}tcPr")
                    tc.insert(0, tcPr)
                # marL/marR/marT/marB en EMU (1pt = 12700 EMU)
                tcPr.set("marL", "91440")   # 7.2pt left
                tcPr.set("marR", "91440")   # 7.2pt right
                tcPr.set("marT", "45720")   # 3.6pt top
                tcPr.set("marB", "45720")   # 3.6pt bottom

        # Solo procesamos la primera tabla encontrada
        break


# ════════════════════════════════════════════════════════════════════════════
# OTRAS DIAPOSITIVAS – Animaciones generales
# ════════════════════════════════════════════════════════════════════════════
def animate_slide1(slide):
    add_fade_transition(slide)
    add_animations(slide, [
        (4,  0,    "fade",       700),
        (5,  300,  "fade",       500),
        (6,  500,  "fly_left",   500),
        (7,  700,  "fade",       500),
        (8,  900,  "fade",       400),
        (9,  1100, "fade",       400),
        (10, 1100, "fade",       400),
        (11, 1100, "fade",       400),
        (12, 1100, "fade",       400),
    ])


def animate_slide2(slide):
    add_fade_transition(slide)
    add_animations(slide, [
        (2, 0,   "fly_bottom", 600),
        (3, 500, "fade",       500),
        (4, 800, "fade",       500),
    ])


def animate_slide4(slide):
    add_fade_transition(slide)
    add_animations(slide, [
        (3, 0,   "fade", 500),
        (4, 300, "fade", 500),
        (5, 500, "fade", 400),
        (6, 800, "fade", 700),
        (7, 1200,"fade", 400),
        (8, 1200,"fade", 400),
    ])


def animate_slide5(slide):
    add_fade_transition(slide)
    add_animations(slide, [
        (2, 0,   "fly_bottom", 600),
        (3, 500, "fade",       500),
        (4, 800, "fade",       500),
        (5, 1050,"fade",       400),
    ])


def animate_slide6(slide):
    add_fade_transition(slide)
    add_animations(slide, [
        (3, 0,    "fade", 500),
        (4, 300,  "fade", 500),
        # Diego
        (5,  700, "fade", 400), (6, 700, "fade", 400), (7, 700, "fade", 400),
        (8,  700, "fade", 400), (9, 700, "fade", 400), (11,700, "fade", 400),
        # María
        (12,1100, "fade", 400), (13,1100,"fade",400), (14,1100,"fade",400),
        (15,1100, "fade", 400), (16,1100,"fade",400), (34,1100,"fade",400),
        # Rosa
        (19,1500, "fade", 400), (20,1500,"fade",400), (21,1500,"fade",400),
        (22,1500, "fade", 400), (23,1500,"fade",400), (35,1500,"fade",400),
        # Ibtisam
        (26,1900, "fade", 400), (27,1900,"fade",400), (28,1900,"fade",400),
        (29,1900, "fade", 400), (30,1900,"fade",400), (33,1900,"fade",400),
        # Leyenda
        (36,2300, "fade", 500),
    ])


def animate_slide7(slide):
    add_fade_transition(slide)
    add_animations(slide, [
        (2,  0,   "fade", 400),
        (4,  0,   "fade", 400),
        (8,  400, "fade", 300), (9, 400,"fade",300), (10,400,"fade",300),
        (11, 400, "fade", 300), (12,400,"fade",300),
        # horas (todos juntos a los 700ms)
        *[(sid, 700, "fade", 300) for sid in [
            14,16,18,20,22,24,26,28,30,32,34,36,38,40,42,44,46,48
        ]],
        # eventos (todos a los 1100ms)
        *[(sid, 1100, "fade", 400) for sid in [
            54,56,58,60,62,64,66,68,70,72,74,76,78,80,82,
            84,86,88,90,92,94,96,98,100,102,104,106,108,110,112
        ]],
        # leyenda
        *[(sid, 1600, "fade", 400) for sid in [114,116,118,120,122,124]],
    ])


def animate_slide8(slide):
    add_fade_transition(slide)
    add_animations(slide, [
        (4,  300,  "fade",       600),
        (12, 1000, "fly_bottom", 600),
        (13, 1500, "fade",       500),
    ])


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════
def enhance(input_path, output_path):
    prs = Presentation(input_path)
    slides = prs.slides

    print("Slide 1: portada + animaciones...")
    animate_slide1(slides[0])

    print("Slide 2: sección FASE 1...")
    animate_slide2(slides[1])

    print("Slide 3: punto de partida + animación ascenso Javier...")
    add_fade_transition(slides[2])
    enhance_slide3_javier(slides[2])

    print("Slide 4: tabla mejorada + animaciones...")
    improve_slide4_table(slides[3])
    animate_slide4(slides[3])

    print("Slide 5: sección FASE 2...")
    animate_slide5(slides[4])

    print("Slide 6: plan persona a persona...")
    animate_slide6(slides[5])

    print("Slide 7: calendario semanal...")
    animate_slide7(slides[6])

    print("Slide 8: cierre...")
    animate_slide8(slides[7])

    prs.save(output_path)
    print(f"\n✓ Presentación guardada en: {output_path}")


if __name__ == "__main__":
    enhance(
        "santa-lucia/Caso2_Presentacion_Javier_2_1.pptx",
        "santa-lucia/Caso2_Presentacion_Javier_2_1_mejorada.pptx",
    )
